from __future__ import annotations

import json
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from app.core.constants import BlockType, SupportedFileType
from app.ingestion.parsers.csv_parser import CsvParser
from app.ingestion.parsers.docx_parser import DocxParser
from app.ingestion.parsers.factory import DEFAULT_PARSER_FACTORY, parse_document
from app.ingestion.parsers.json_parser import JsonParser
from app.ingestion.parsers.pdf_parser import PdfParser
from app.ingestion.parsers.pptx_parser import PptxParser
from app.ingestion.parsers.xlsx_parser import XlsxParser


def test_csv_parser_creates_table_blocks(tmp_path):
    file_path = tmp_path / "scores.csv"
    file_path.write_text("name,score\nAda,98\nLinus,95\n", encoding="utf-8")

    document = CsvParser().parse(file_path)

    assert document.file_type == SupportedFileType.CSV.value
    assert [block.block_type for block in document.blocks] == [
        BlockType.TABLE.value,
        BlockType.TABLE.value,
    ]
    assert document.blocks[0].text == "name: Ada | score: 98"
    assert document.blocks[0].source_location.line_start == 2
    assert document.blocks[0].metadata["columns"] == ["name", "score"]


def test_json_parser_creates_json_path_blocks(tmp_path):
    file_path = tmp_path / "profile.json"
    file_path.write_text(
        json.dumps(
            {
                "team": "search",
                "members": [
                    {"name": "Ada", "role": "engineer"},
                    {"name": "Linus", "role": "reviewer"},
                ],
            }
        ),
        encoding="utf-8",
    )

    document = JsonParser().parse(file_path)

    assert document.file_type == SupportedFileType.JSON.value
    assert [block.block_type for block in document.blocks] == [
        BlockType.JSON.value,
        BlockType.JSON.value,
        BlockType.JSON.value,
    ]
    assert document.blocks[0].metadata["json_path"] == "$.team"
    assert document.blocks[1].metadata["json_path"] == "$.members[0]"
    assert '"name": "Ada"' in document.blocks[1].text


def test_docx_parser_preserves_paragraphs_lists_and_tables(tmp_path):
    file_path = tmp_path / "handbook.docx"
    document = Document()
    document.core_properties.title = "Handbook"
    document.add_heading("Handbook", level=1)
    document.add_paragraph("Welcome to the handbook.")
    document.add_paragraph("Bring your laptop.", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Policy"
    table.cell(0, 1).text = "Owner"
    table.cell(1, 0).text = "Leave"
    table.cell(1, 1).text = "People"
    document.save(file_path)

    parsed_document = DocxParser().parse(file_path)

    assert parsed_document.title == "Handbook"
    assert parsed_document.file_type == SupportedFileType.DOCX.value
    assert [block.block_type for block in parsed_document.blocks] == [
        BlockType.HEADING.value,
        BlockType.PARAGRAPH.value,
        BlockType.LIST_ITEM.value,
        BlockType.TABLE.value,
    ]
    assert parsed_document.blocks[3].metadata["format"] == "docx_table"
    assert "Leave" in parsed_document.blocks[3].text


def test_pptx_parser_preserves_slide_titles_text_and_tables(tmp_path):
    file_path = tmp_path / "roadmap.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap"
    slide.placeholders[1].text = "Launch search beta"
    table_shape = slide.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(3),
        Inches(5),
        Inches(1),
    )
    table_shape.table.cell(0, 0).text = "Milestone"
    table_shape.table.cell(0, 1).text = "Status"
    table_shape.table.cell(1, 0).text = "Parser"
    table_shape.table.cell(1, 1).text = "Ready"
    presentation.save(file_path)

    document = PptxParser().parse(file_path)

    assert document.title == "Roadmap"
    assert document.file_type == SupportedFileType.PPTX.value
    assert [block.block_type for block in document.blocks] == [
        BlockType.HEADING.value,
        BlockType.PARAGRAPH.value,
        BlockType.TABLE.value,
    ]
    assert document.blocks[0].source_location.slide_number == 1
    assert "Parser" in document.blocks[2].text


def test_xlsx_parser_creates_sheet_heading_and_row_blocks(tmp_path):
    file_path = tmp_path / "metrics.xlsx"
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Metrics"
    worksheet.append(["name", "value"])
    worksheet.append(["latency", 120])
    worksheet.append(["quality", "high"])
    workbook.save(file_path)

    document = XlsxParser().parse(file_path)

    assert document.file_type == SupportedFileType.XLSX.value
    assert [block.block_type for block in document.blocks] == [
        BlockType.HEADING.value,
        BlockType.TABLE.value,
        BlockType.TABLE.value,
    ]
    assert document.blocks[0].text == "Metrics"
    assert document.blocks[1].text == "name: latency | value: 120"
    assert document.blocks[1].source_location.sheet_name == "Metrics"
    assert document.blocks[1].source_location.row_start == 2


def test_pdf_parser_extracts_text_blocks(tmp_path):
    file_path = tmp_path / "sample.pdf"
    _write_text_pdf(file_path, "Hello PDF parser")

    document = PdfParser().parse(file_path)

    assert document.file_type == SupportedFileType.PDF.value
    assert document.blocks[0].block_type == BlockType.PARAGRAPH.value
    assert "Hello PDF parser" in document.blocks[0].text
    assert document.blocks[0].source_location.page_number == 1


def test_default_factory_supports_all_declared_parser_formats(tmp_path):
    supported_extensions = DEFAULT_PARSER_FACTORY.supported_extensions

    assert {
        ".txt",
        ".md",
        ".markdown",
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".csv",
        ".json",
    }.issubset(supported_extensions)

    file_path = tmp_path / "factory.json"
    file_path.write_text('{"status": "ok"}', encoding="utf-8")

    document = parse_document(file_path)

    assert document.file_type == SupportedFileType.JSON.value


def _write_text_pdf(file_path: Path, text: str) -> None:
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)

    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_reference = writer._add_object(font)

    resources = DictionaryObject(
        {
            NameObject("/Font"): DictionaryObject(
                {NameObject("/F1"): font_reference}
            )
        }
    )
    page[NameObject("/Resources")] = resources

    stream = DecodedStreamObject()
    escaped_text = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream.set_data(f"BT /F1 24 Tf 100 700 Td ({escaped_text}) Tj ET".encode("utf-8"))
    page[NameObject("/Contents")] = writer._add_object(stream)

    with file_path.open("wb") as pdf_file:
        writer.write(pdf_file)
