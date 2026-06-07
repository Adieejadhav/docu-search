from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from app.core.constants import BlockType, SupportedFileType, SUPPORTED_PPTX_EXTENSIONS
from app.core.exceptions import ParserError
from app.ingestion.normalizers.block_schema import (
    DocumentBlock,
    ParsedDocument,
    SourceLocation,
)
from app.ingestion.parsers._tabular import format_markdown_table, normalize_row
from app.ingestion.parsers.base import BaseDocumentParser


class PptxParser(BaseDocumentParser):
    """
    Parser for PowerPoint .pptx presentations.

    Slide titles become heading blocks; text frame paragraphs and tables become
    paragraph/list/table blocks with slide source locations.
    """

    supported_extensions = SUPPORTED_PPTX_EXTENSIONS

    def parse(self, file_path: str | Path) -> ParsedDocument:
        path = self.validate_input_file(file_path)
        presentation = self._open_presentation(path)
        blocks = self._parse_presentation(presentation)
        title = self._extract_title(presentation, blocks) or self.get_document_title(path)

        return self.build_document(
            path=path,
            title=title,
            file_type=SupportedFileType.PPTX,
            blocks=blocks,
            metadata={
                "slide_count": len(presentation.slides),
                "core_properties_title": presentation.core_properties.title or None,
            },
        )

    def _open_presentation(self, path: Path) -> Presentation:
        try:
            return Presentation(str(path))
        except PackageNotFoundError as exc:
            raise ParserError(
                f"Unable to open PPTX file: {path}",
                code="PPTX_OPEN_FAILED",
                details={"file_path": str(path)},
            ) from exc
        except Exception as exc:
            raise ParserError(
                f"Unable to parse PPTX file: {path}",
                code="PPTX_PARSE_FAILED",
                details={"file_path": str(path)},
            ) from exc

    def _parse_presentation(self, presentation: Presentation) -> list[DocumentBlock]:
        blocks: list[DocumentBlock] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_title = self._slide_title(slide)

            if slide_title:
                blocks.append(
                    DocumentBlock(
                        block_type=BlockType.HEADING,
                        text=slide_title,
                        order=len(blocks),
                        level=1,
                        source_location=SourceLocation(slide_number=slide_number),
                        metadata={"slide_number": slide_number, "role": "slide_title"},
                    )
                )

            for shape in slide.shapes:
                if slide.shapes.title is not None and shape == slide.shapes.title:
                    continue

                if getattr(shape, "has_table", False):
                    table_text = self._shape_table_to_text(shape)
                    if table_text:
                        blocks.append(
                            DocumentBlock(
                                block_type=BlockType.TABLE,
                                text=table_text,
                                order=len(blocks),
                                parent_path=[slide_title] if slide_title else [],
                                source_location=SourceLocation(
                                    slide_number=slide_number
                                ),
                                metadata={
                                    "format": "pptx_table",
                                    "slide_number": slide_number,
                                    "shape_name": shape.name,
                                },
                            )
                        )
                    continue

                if getattr(shape, "has_text_frame", False):
                    blocks.extend(
                        self._text_frame_blocks(
                            shape,
                            slide_number=slide_number,
                            parent_path=[slide_title] if slide_title else [],
                            start_order=len(blocks),
                        )
                    )

        return blocks

    def _text_frame_blocks(
        self,
        shape: Any,
        *,
        slide_number: int,
        parent_path: list[str],
        start_order: int,
    ) -> list[DocumentBlock]:
        blocks: list[DocumentBlock] = []

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            block_type = (
                BlockType.LIST_ITEM if paragraph.level and paragraph.level > 0 else BlockType.PARAGRAPH
            )
            blocks.append(
                DocumentBlock(
                    block_type=block_type,
                    text=text,
                    order=start_order + len(blocks),
                    parent_path=parent_path,
                    source_location=SourceLocation(slide_number=slide_number),
                    metadata={
                        "slide_number": slide_number,
                        "shape_name": shape.name,
                        "paragraph_level": paragraph.level,
                    },
                )
            )

        return blocks

    def _shape_table_to_text(self, shape: Any) -> str:
        rows = [
            normalize_row(cell.text for cell in row.cells)
            for row in shape.table.rows
        ]
        return format_markdown_table(rows)

    def _slide_title(self, slide: Any) -> str | None:
        title_shape = slide.shapes.title
        if title_shape is None or not getattr(title_shape, "has_text_frame", False):
            return None

        title = title_shape.text.strip()
        return title or None

    def _extract_title(
        self,
        presentation: Presentation,
        blocks: list[DocumentBlock],
    ) -> str | None:
        core_title = presentation.core_properties.title
        if core_title and core_title.strip():
            return core_title.strip()

        for block in blocks:
            if block.block_type == BlockType.HEADING and block.level == 1:
                return block.text

        return None
